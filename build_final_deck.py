"""
Final Project — Demo & Presentation deck, on the Booth 16:9 master.

    python3 build_final_deck.py

Ten minutes, five speakers. Every figure is from score_rubric.py /
stability_test.py; every problem-side figure is from Group Assignment 1.
Speaker notes carry the spoken script, so the deck and the speech stay in sync.

The Booth canvas is 26.67 x 15 in, roughly 2x a normal 13.3in deck, so type
sizes here are about double what they would be elsewhere.
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from speaker_notes import SCRIPT, TIMING

HERE = Path(__file__).parent
TEMPLATE = HERE / "booth_template.pptx"
OUT = HERE.parent / "Final_Presentation_Energy_Anomaly_Explainer.pptx"

BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x53, 0x53, 0x53)
LGREY = RGBColor(0xD8, 0xD8, 0xD8)
GOLD = RGBColor(0xF5, 0xC2, 0x01)
BLUE = RGBColor(0x52, 0x6D, 0xB0)
ORANGE = RGBColor(0xDC, 0x59, 0x24)
GREEN = RGBColor(0x1A, 0x7F, 0x37)
RED = RGBColor(0xB3, 0x26, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xF4, 0xF4, 0xF4)

W, H = 26.67, 15.0
L, R = 0.95, 0.95
CW = W - L - R          # 24.77 usable width
FONT = "Arial"

prs = Presentation(str(TEMPLATE))

# strip the template's seven sample slides
ids = prs.slides._sldIdLst
for sid in list(ids):
    prs.part.drop_rel(sid.rId)
    ids.remove(sid)


def strip_slide_numbers(prs):
    """Remove the page numbers the Booth template paints on every slide.

    They come from two different mechanisms and both have to go. The master and
    layouts carry sldNum placeholders holding a <a:fld type="slidenum"> field;
    separately, layouts 3-6 carry an ordinary text shape with the literal string
    "01" sitting in the footer band. Only the second one is what actually shows
    on a Blank-layout slide, and it is invisible to any search for a field."""
    ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    removed = 0
    for part in [prs.slide_master, *prs.slide_layouts]:
        for shp in list(part.shapes):
            ph = shp._element.find(f"{ns}nvSpPr/{ns}nvPr/{ns}ph")
            is_field = ph is not None and ph.get("type") == "sldNum"
            is_literal = (shp.has_text_frame
                          and shp.text_frame.text.strip().isdigit())
            if is_field or is_literal:
                shp._element.getparent().remove(shp._element)
                removed += 1
    return removed


strip_slide_numbers(prs)

LY_TITLE, LY_TITLECONTENT, LY_BLANK = 0, 2, 6


# ── primitives ───────────────────────────────────────────────────────────────
def add(layout=LY_BLANK):
    return prs.slides.add_slide(prs.slide_layouts[layout])


def drop_empty_placeholders(slide):
    """Every slide's content is in explicit text boxes, so no placeholder is
    load-bearing — including the layout's slide-number one, which we do not want."""
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def tb(slide, x, y, w, h, parts, size=26, color=BLACK, bold=False, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, space=6, line=0.92):
    """parts: str | [(text, {opts})] | [[run,...], [run,...]] for multiple paragraphs."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    if isinstance(parts, str):
        paras = [[(parts, {})]]
    elif parts and isinstance(parts[0], tuple):
        paras = [parts]
    else:
        paras = parts

    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(space)
        if isinstance(runs, str):
            runs = [(runs, {})]
        for text, o in runs:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = FONT
            f.size = Pt(o.get("size", size))
            f.bold = o.get("b", bold)
            f.italic = o.get("i", False)
            f.color.rgb = o.get("c", color)
    return box


def title(slide, text, sub=None):
    tb(slide, L, 0.72, CW, 1.5, text, size=60, bold=True)
    if sub:
        tb(slide, L, 2.32, CW, 0.9, sub, size=27, color=GREY)
    return 3.45 if sub else 2.60


def card(slide, x, y, w, h, fill=TINT, line=None, lw=1.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.text_frame.text = ""
    return s


def stat(slide, x, y, w, big, label, note=None, c=BLACK, bigsize=76):
    tb(slide, x, y, w, 1.5, big, size=bigsize, bold=True, color=c, line=0.85)
    tb(slide, x, y + 1.42, w, 0.95, label, size=24, color=BLACK, line=1.0)
    if note:
        tb(slide, x, y + 2.42, w, 1.2, note, size=19, color=GREY, line=1.02)


def bullets(slide, x, y, w, items, size=26, gap=10, colour=BLACK):
    paras = []
    for it in items:
        runs = [(it, {})] if isinstance(it, str) else it
        paras.append([("•   ", {"c": GOLD, "b": True})] + list(runs))
    return tb(slide, x, y, w, 1.7 * len(items), paras, size=size, color=colour, space=gap)


def table(slide, x, y, w, cols, rows, widths, sizes=(21, 20), head_fill=BLACK,
          row_h=0.62, head_h=0.62):
    n_r, n_c = len(rows) + 1, len(cols)
    g = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w),
                               Inches(head_h + row_h * len(rows))).table
    g.first_row = True
    g.horz_banding = False
    for i, cw in enumerate(widths):
        g.columns[i].width = Inches(cw)
    g.rows[0].height = Inches(head_h)
    for i in range(1, n_r):
        g.rows[i].height = Inches(row_h)

    def put(cell, content, size, bold, fill, color=BLACK, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.margin_left = cell.margin_right = Inches(0.14)
        cell.margin_top = cell.margin_bottom = Inches(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        p.line_spacing = 0.92
        for text, o in ([(content, {})] if isinstance(content, str) else content):
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = FONT
            f.size = Pt(o.get("size", size))
            f.bold = o.get("b", bold)
            f.italic = o.get("i", False)
            f.color.rgb = o.get("c", color)

    for i, c in enumerate(cols):
        put(g.cell(0, i), c, sizes[0], True, head_fill, WHITE,
            PP_ALIGN.CENTER if i else PP_ALIGN.LEFT)
    for ri, row in enumerate(rows):
        fill = WHITE if ri % 2 == 0 else TINT
        for ci, c in enumerate(row):
            put(g.cell(ri + 1, ci), c, sizes[1], False, fill,
                align=PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT)
    return g


def stamp(slide, n):
    """Speaker notes and the corner credit, both from speaker_notes.py."""
    slide.notes_slide.notes_text_frame.text = SCRIPT[n]
    who, secs = TIMING[n]
    tb(slide, W - R - 7.0, H - 1.05, 7.0, 0.5,
       f"{who}   ·   {secs // 60}:{secs % 60:02d}",
       size=18, color=GREY, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# 1 — Title
# ═════════════════════════════════════════════════════════════════════════════
s = add(LY_TITLE)
drop_empty_placeholders(s)
tb(s, L, 4.6, 22.0, 2.4, "Energy Anomaly Explainer", size=96, bold=True)
tb(s, L, 7.35, 22.0, 1.2,
   "Should I send a technician? An answer in ten seconds, not three days.",
   size=34, color=GREY)
card(s, L, 9.15, 3.2, 0.06, fill=GOLD)
tb(s, L, 9.85, 24.0, 1.6,
   [[("AI Analytics Aces", {"b": True, "size": 26})],
    [("Kathleen (Helen) Baca · Alfreda Holloway · Nathanael Martin-Nelson · "
      "Cindy Rozental · Yun Zhao", {"size": 22, "c": GREY})],
    [("BUSN 43800 · Analytics with AI · Chicago Booth", {"size": 22, "c": GREY})]],
   space=7)
stamp(s, 1)

# ═════════════════════════════════════════════════════════════════════════════
# 2 — The problem
# ═════════════════════════════════════════════════════════════════════════════
s = add()
y = title(s, "The problem, as a measurable gap",
          "Operations managers at U.S. distribution centers of 100,000 sq ft or more")
cw, gap = 5.75, 0.59
for i, (big, lab, note, c) in enumerate([
    ("1–3 days", "to diagnose one spike today",
     "Requires a facilities engineer, at $500–$2,000 per incident", BLACK),
    ("3–5 / month", "spikes per facility that trigger manual investigation",
     "30–50 open events a month across a 10-facility network", BLACK),
    ("167 : 1", "cost of a missed fault vs. a needless dispatch",
     "A dispatch is ~$300. A missed compressor fault runs to five figures", ORANGE),
    ("$60–240k", "avoidable cost a year, 10-facility network",
     "Excess consumption and spoilage from delayed diagnosis alone", BLACK),
]):
    x = L + i * (cw + gap)
    card(s, x, y, cw, 5.4, fill=TINT)
    stat(s, x + 0.42, y + 0.5, cw - 0.84, big, lab, note, c=c, bigsize=64)

yy = y + 6.0
card(s, L, yy, CW, 1.65, fill=None, line=BLACK, lw=2)
tb(s, L + 0.5, yy + 0.34, CW - 1.0, 1.0,
   [("The decision that turns on it.  ", {"b": True}),
    ("Dispatch a technician with a described symptom, monitor for recurrence, or "
     "dismiss it as a known operational event — inside the same shift, without "
     "calling an engineer.", {})], size=27)
stamp(s, 2)

# ═════════════════════════════════════════════════════════════════════════════
# 3 — Two theories
# ═════════════════════════════════════════════════════════════════════════════
s = add()
y = title(s, "We committed to a test before we built anything",
          "Two rival accounts of why these spikes happen — only one leaves room for our product")
cw = (CW - 0.7) / 2
card(s, L, y, cw, 5.0, fill=TINT)
tb(s, L + 0.5, y + 0.42, cw - 1.0, 0.6, "THEORY A", size=22, bold=True, color=BLUE)
tb(s, L + 0.5, y + 1.15, cw - 1.0, 3.4,
   "Spikes carry diagnostic signal. Timing against shift boundaries, magnitude "
   "against the facility's own history, and duration against equipment cycle "
   "times are jointly enough to classify cause.", size=27)
tb(s, L + 0.5, y + 3.95, cw - 1.0, 0.7,
   [("→  The product works.", {"b": True, "c": BLUE})], size=25)

x2 = L + cw + 0.7
card(s, x2, y, cw, 5.0, fill=TINT)
tb(s, x2 + 0.5, y + 0.42, cw - 1.0, 0.6, "THEORY B", size=22, bold=True, color=ORANGE)
tb(s, x2 + 0.5, y + 1.15, cw - 1.0, 3.4,
   "Operational variation is structurally indistinguishable from equipment "
   "faults on kWh alone. Busy days, overtime, weather, rented equipment — the "
   "features carry noise, not signal.", size=27)
tb(s, x2 + 0.5, y + 3.95, cw - 1.0, 0.7,
   [("→  The product degrades the decision it replaces.", {"b": True, "c": ORANGE})],
   size=25)

yy = y + 5.6
card(s, L, yy, CW, 2.05, fill=None, line=BLACK, lw=2)
tb(s, L + 0.5, yy + 0.38, CW - 1.0, 1.7,
   [[("The evidence we said would settle it, before we saw any results.  ", {"b": True}),
     ("Precision ≥ 75% and recall ≥ 70% on the equipment fault class, against "
      "held-out labels.", {})],
    [("Everything that follows is us reporting against that bar.", {"i": True, "c": GREY})]],
   size=27, space=9)
stamp(s, 3)

# ═════════════════════════════════════════════════════════════════════════════
# 4 — The product
# ═════════════════════════════════════════════════════════════════════════════
s = add()
y = title(s, "What we built",
          "Hourly meter data in, a decision the manager can act on out — no new hardware")

bw, bgap = 4.55, 0.62
steps = [
    ("Hourly kWh", "Smart meter or BMS export, uploaded as CSV", GREY),
    ("Real weather", "Open-Meteo hourly temperature, fetched automatically", GREY),
    ("Detector", "Flags spikes against a rolling baseline — deliberately weather-blind, "
                 "like the alarm a manager already gets", BLACK),
    ("Agent, 4 tools", "Fetches readings, builds a weather-matched baseline, runs a "
                       "significance test, finds comparable past events", BLACK),
    ("Claude Opus 5", "Classifies into 14 causes, writes the explanation and the action, "
                      "schema-enforced", BLUE),
]
for i, (h, d, c) in enumerate(steps):
    x = L + i * (bw + bgap)
    card(s, x, y, bw, 3.9, fill=TINT)
    tb(s, x + 0.36, y + 0.38, bw - 0.72, 0.8, h, size=27, bold=True, color=c)
    tb(s, x + 0.36, y + 1.32, bw - 0.72, 2.3, d, size=21, color=GREY, line=1.05)
    if i < len(steps) - 1:
        tb(s, x + bw + 0.06, y + 1.32, 0.5, 0.8, "›", size=40, color=LGREY,
           align=PP_ALIGN.CENTER)

yy = y + 4.5
card(s, L, yy, CW, 3.25, fill=None, line=BLACK, lw=2)
tb(s, L + 0.55, yy + 0.4, CW - 1.1, 0.7,
   "Seven fields come back, and a manager can act on all of them", size=29, bold=True)
tb(s, L + 0.55, yy + 1.25, CW - 1.1, 1.7,
   [[("Cause and subtype · confidence · a plain-English explanation · ", {}),
     ("dispatch, monitor, or dismiss", {"b": True}),
     (" · the manager's next step · the symptom the technician should look for.", {})],
    [("A fifth path is always open: record an ", {"c": GREY}),
     ("exception", {"i": True, "c": GREY}),
     (" in free text when none of the three fit — that is how the product tells "
      "us what its own taxonomy is missing.", {"c": GREY})]],
   size=25, space=10)
stamp(s, 4)

# ═════════════════════════════════════════════════════════════════════════════
# 5 — Live demo
# ═════════════════════════════════════════════════════════════════════════════
s = add()
card(s, 0, 0, W, H, fill=BLACK)
tb(s, L, 1.5, CW, 1.0, "DEMO — SCREEN RECORDING", size=24, bold=True, color=GOLD)
tb(s, L, 2.6, 23.0, 2.2, "It works, it saves a trip, and it fails",
   size=64, bold=True, color=WHITE)

items = [
    ("ANO-2000 — it works.", "  Refrigeration Zone A at 4am, 122.5 kWh against a 61.1 "
     "baseline — double, for three hours, at 48°F. Compressor fault, 0.78, dispatch, "
     "with a symptom a refrigeration tech can act on."),
    ("ANO-2012 — it saves the trip.", "  Lighting Grid South at 3.75× normal, the "
     "scariest number in the set. The hour before recorded almost nothing: the meter "
     "dropped out and caught up. Dismiss at 0.88 — no technician sent."),
    ("ANO-2009 — it fails, and we show you.", "  Eight hours at 1.15× on the Compressor "
     "Bank, called a fault at 0.78. It was a Peak Throughput Day, and no prompt fixes "
     "that when shift schedules are not an input."),
    ("Then the guardrails.", "  A meter file in watts instead of kilowatts, caught on "
     "upload. And an exception recorded when none of the three actions fit."),
]
yy = 5.6
for head, rest in items:
    card(s, L, yy + 0.28, 0.13, 0.13, fill=GOLD)
    tb(s, L + 0.55, yy, CW - 0.55, 1.5,
       [(head, {"b": True, "c": WHITE}), (rest, {"c": RGBColor(0xC8, 0xC8, 0xC8)})],
       size=26, line=1.06)
    yy += 1.72

tb(s, L, H - 1.6, CW, 0.6,
   "Recorded on the presentation machine, unedited except for cuts between cases "
   "— we are telling you it is a recording, not a live run",
   size=20, color=RGBColor(0x9A, 0x9A, 0x9A))
stamp(s, 5)

# ═════════════════════════════════════════════════════════════════════════════
# 6 — The rubric
# ═════════════════════════════════════════════════════════════════════════════
s = add()
y = title(s, "How we decide an output is good",
          "One output = one classified spike. Five criteria, scored on every case, "
          "each with a threshold fixed in advance")
table(s, L, y, CW,
      ["Criterion", "How it is judged", "Acceptable at"],
      [["Correctness", "Class and subtype vs. the held-out label, exact match", "Both exact"],
       ["Completeness", "All seven fields present; a symptom whenever the action is dispatch",
        "All fields"],
       ["Calibration", "Commit (≥0.75) when right, hedge when wrong. A wrong label at ≥0.75 "
                       "is an overclaim", "No overclaims"],
       ["Usability", "Names a cause · states evidence · specific next step · no jargon. "
                     "Two blind human raters", "≥ 80% of cases"],
       ["Timeliness", "Wall clock from request to returned record", "< 5 minutes"]],
      [4.6, 15.0, 5.17], sizes=(23, 23), row_h=1.06, head_h=0.75)

yy = y + 6.4
card(s, L, yy, CW, 2.3, fill=TINT)
tb(s, L + 0.55, yy + 0.35, CW - 1.1, 1.85,
   [[("Four of the five need no judge at all.  ", {"b": True}),
     ("Correctness is a string match against labels the classifier never sees. "
      "Completeness and Calibration are field and threshold checks. Timeliness "
      "reads a timer. The model is never both author and marker.", {})],
    [("Usability is the exception — so two of us score it blind, without the label "
      "or the confidence. An LLM pre-rating agrees on 13 of 14, and we do not quote "
      "that as the result until the human scores land.", {"c": GREY})]],
   size=24, space=9)
stamp(s, 6)

# ═════════════════════════════════════════════════════════════════════════════
# 7 — Results
# ═════════════════════════════════════════════════════════════════════════════
s = add()
y = title(s, "What the test set says — including the failures",
          "14 held-out cases: 9 equipment faults, 3 operational variations, 2 data anomalies")

cw3 = (CW - 1.2) / 3
for i, (big, lab, c) in enumerate([
    ("82%", "Precision on equipment fault      bar ≥ 75%", GREEN),
    ("100%", "Recall on equipment fault      bar ≥ 70%", GREEN),
    ("10.0 s", "Average time to classify      bar < 5 min", GREEN),
]):
    x = L + i * (cw3 + 0.6)
    card(s, x, y, cw3, 2.5, fill=TINT)
    tb(s, x + 0.45, y + 0.35, cw3 - 0.9, 1.0, big, size=62, bold=True, color=c)
    tb(s, x + 0.45, y + 1.62, cw3 - 0.9, 0.7, lab, size=21, color=GREY)

yy = y + 3.05
tb(s, L, yy, CW, 0.7,
   [("Theory A cleared its bar. And three criteria still failed.", {"b": True})], size=31)

table(s, L, yy + 0.9, CW,
      ["Failed criterion", "Result", "What went wrong", "Diagnosis"],
      [[[("Correctness", {"b": True})], [("86% class / 50% subtype", {"c": RED, "b": True})],
        "Peak Throughput Day and Temporary Equipment Rental both read as faults. "
        "Neither shift schedules nor rental logs are inputs we receive",
        [("Specification", {"b": True})]],
       [[("Calibration", {"b": True})], [("1 overclaim, 43%", {"c": RED, "b": True})],
        "Confidence averages 0.70 when right and 0.75 when wrong — inverted, "
        "not merely noisy", [("Model", {"b": True})]],
       [[("Decision value", {"b": True})], [("$35,200 vs $7,500", {"c": RED, "b": True})],
        "Accurate labels, and following the tool still costs 4.7× dispatching on "
        "everything", [("Question", {"b": True})]]],
      [4.3, 4.5, 12.47, 3.5], sizes=(22, 21), row_h=1.28, head_h=0.72)
stamp(s, 7)

# ═════════════════════════════════════════════════════════════════════════════
# 8 — Decision value
# ═════════════════════════════════════════════════════════════════════════════
s = add()
y = title(s, "The failure that matters: an accurate tool that loses money",
          "We priced the tool against two fixed policies — $300 a dispatch, $2,000 a "
          "missed fault, over all 25 anomalies")
table(s, L, y, CW,
      ["Policy", "Sent", "Faults caught", "Dispatch spend", "Missed-fault exposure", "Total cost"],
      [[[("Follow the tool", {"b": True})], "4", "3 / 20", "$1,200", "17 × $2,000 = $34,000",
        [("$35,200", {"b": True, "c": RED})]],
       ["Dispatch on every spike", "25", "20 / 20", "$7,500", "none",
        [("$7,500", {"b": True, "c": GREEN})]],
       ["Dispatch on none", "0", "0 / 20", "$0", "20 × $2,000 = $40,000",
        [("$40,000", {"b": True, "c": RED})]]],
      [5.4, 2.3, 3.3, 4.0, 5.87, 3.9], sizes=(22, 23), row_h=0.95, head_h=0.72)

yy = y + 4.35
cwh = (CW - 0.7) / 2
card(s, L, yy, cwh, 3.7, fill=TINT)
tb(s, L + 0.5, yy + 0.38, cwh - 1.0, 0.7, "97% of the cost is exposure, not spend",
   size=29, bold=True)
tb(s, L + 0.5, yy + 1.3, cwh - 1.0, 2.2,
   "The tool spends $1,200 and leaves $34,000 of faults uninvestigated. Only 4 of "
   "25 spikes clear the 0.75 confidence gate, so 17 of 20 real faults are told to "
   "“monitor” and nobody is sent.", size=25, line=1.05)

x2 = L + cwh + 0.7
card(s, x2, yy, cwh, 3.7, fill=None, line=BLACK, lw=2)
tb(s, x2 + 0.5, yy + 0.38, cwh - 1.0, 0.7, "We asked the wrong question",
   size=29, bold=True)
tb(s, x2 + 0.5, yy + 1.3, cwh - 1.0, 2.2,
   "We built for “what is this spike?”. The decision needs “should I dispatch?”. "
   "A correct label at 0.72 confidence still sends nobody — so accuracy and "
   "usefulness came apart.", size=25, line=1.05)
stamp(s, 8)

# ═════════════════════════════════════════════════════════════════════════════
# 9 — Limitations
# ═════════════════════════════════════════════════════════════════════════════
s = add()
y = title(s, "What it cannot do",
          "From our two independent validation audits — stated here rather than "
          "buried in an appendix")
rows = [
    ("It cannot tell a busy day from a broken compressor.",
     "Shift schedules, throughput and rental logs all move load, and none is an input we "
     "receive. This is Theory B, and it is real."),
    ("We tried the fix our own spec prescribed. It failed.",
     "Adding shift schedules dropped precision from 82% to 64%: a day planned at 130% adds "
     "7–10 kWh, the alarm fires at 20, so a realistic busy day never reaches the detector."),
    ("Confidence should not be trusted as a probability.",
     "It is inverted — 0.70 when right, 0.75 when wrong. Do not build a dispatch gate on it, "
     "which is exactly the mistake we made."),
    ("Every result is on synthetic data.",
     "Real temperature from Open-Meteo, but generated consumption. No maintenance logs, "
     "so no validation against a fault someone actually found."),
    ("The classifier never faces a false alarm.",
     "The detector raises 562 alarms a year; only the 25 matching an injected anomaly are "
     "scored. Precision measures label accuracy, never real-event-vs-false-alarm."),
]
yy = y
for head, rest in rows:
    card(s, L, yy, CW, 1.42, fill=TINT)
    tb(s, L + 0.5, yy + 0.22, CW - 1.0, 1.0,
       [(head + "  ", {"b": True}), (rest, {"c": GREY})], size=24, line=1.04)
    yy += 1.58
stamp(s, 9)

# ═════════════════════════════════════════════════════════════════════════════
# 10 — What happens next
# ═════════════════════════════════════════════════════════════════════════════
s = add()
y = title(s, "What happens next",
          "How we would monitor it, what would tell us it has stopped working, "
          "and what we would build first")

cw2 = (CW - 1.2) / 3
blocks = [
    ("Monitor", BLUE, [
        "Exception notes — the free-text reasons managers give when none of the three "
        "actions fit. A reason that recurs is a missing 15th cause",
        "Dispatch outcome: did the technician find the fault we named?",
        "Share of spikes clearing the confidence gate, tracked weekly",
        "Input-guard findings per upload — a rising count means the feed changed",
    ]),
    ("Stop-the-line triggers", ORANGE, [
        "Precision on equipment fault below 60% over a rolling 30 days",
        "Two consecutive weeks where confirmed faults were labelled operational",
        "Exception rate above 20% of classified spikes",
        "Any month where following the tool costs more than dispatching on everything",
    ]),
    ("Build next, in order", GREEN, [
        "Replace the fixed 0.75 gate with an expected-cost rule using the 167:1 ratio — "
        "this is the fix that turns 82% precision into money saved",
        "Move detection from the sub-system to the facility, so site-wide events stop "
        "looking like nine independent faults",
        "Field pilot against real maintenance logs",
    ]),
]
for i, (head, c, items) in enumerate(blocks):
    x = L + i * (cw2 + 0.6)
    card(s, x, y, cw2, 7.35, fill=TINT)
    tb(s, x + 0.45, y + 0.4, cw2 - 0.9, 0.8, head, size=30, bold=True, color=c)
    bullets(s, x + 0.45, y + 1.45, cw2 - 0.9, items, size=22, gap=13)
stamp(s, 10)

# ═════════════════════════════════════════════════════════════════════════════
# 11 — Close
# ═════════════════════════════════════════════════════════════════════════════
s = add()
card(s, 0, 0, W, H, fill=BLACK)
tb(s, L, 2.5, 23.0, 1.0, "WHERE WE LANDED", size=24, bold=True, color=GOLD)
tb(s, L, 3.7, 23.5, 3.2,
   "Theory A cleared the bar we set.\nTheory B was right about where the confusion lives.",
   size=52, bold=True, color=WHITE, line=1.1)

yy = 8.0
cw2 = (CW - 0.7) / 2
tb(s, L, yy, cw2, 3.0,
   [[("The product works as a classifier.", {"b": True, "c": GOLD, "size": 30})],
    [("82% precision, 100% recall, ten seconds a case, and an explanation a "
      "manager can act on without calling anyone.", {"c": RGBColor(0xD5, 0xD5, 0xD5),
                                                     "size": 25})]],
   space=10)
tb(s, L + cw2 + 0.7, yy, cw2, 3.0,
   [[("It is not yet worth deploying.", {"b": True, "c": GOLD, "size": 30})],
    [("The dispatch rule on top of it loses money, and we can show you the "
      "arithmetic. That is the next build, and it is a rule change.",
      {"c": RGBColor(0xD5, 0xD5, 0xD5), "size": 25})]],
   space=10)

card(s, L, 11.9, 3.2, 0.06, fill=GOLD)
tb(s, L, 12.5, 24.0, 1.4,
   [[("Repository, product brief and slides submitted on Canvas · "
      "github.com/yzhao24/ai-analytics-aces", {"size": 22, "c": RGBColor(0x9A, 0x9A, 0x9A)})],
    [("Questions — especially about the failures.",
      {"size": 26, "b": True, "c": WHITE})]],
   space=9)
stamp(s, 11)

# ═════════════════════════════════════════════════════════════════════════════
for sl in prs.slides:
    drop_empty_placeholders(sl)
prs.save(str(OUT))
print("wrote", OUT, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
